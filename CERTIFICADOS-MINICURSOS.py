import pandas as pd
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import re
import os

# Meses em PT-BR
meses = {
    1: "janeiro", 2: "fevereiro", 3: "março", 4: "abril",
    5: "maio", 6: "junho", 7: "julho", 8: "agosto",
    9: "setembro", 10: "outubro", 11: "novembro", 12: "dezembro"
}

# Carregar lista de participantes
df = pd.read_excel("participantes.xlsx")
df["Data"] = pd.to_datetime(df["Data"])
df["DIA"] = df["Data"].dt.day.astype(str)
df["MES"] = df["Data"].dt.month.map(meses)
df["ANO"] = df["Data"].dt.year.astype(str)

# Modelo
TEMPLATE = "Template capacitações.pptx"

# Placeholders que ficam em azul
placeholders_azuis = {
    "{NOME}": {"valor": "Nome", "size": Pt(16)},
    "{CAPACITACAO}": {"valor": "Capacitacao", "size": Pt(16)},
    "{EVENTO}": {"valor": "Evento", "size": Pt(16)},
}

# Criar pasta de saída
os.makedirs("certificados", exist_ok=True)

for _, row in df.iterrows():
    prs = Presentation(TEMPLATE)  # novo arquivo por aluno
    modelo = prs.slides[0]
    slide = prs.slides.add_slide(modelo.slide_layout)

    for shape in modelo.shapes:
        if not shape.has_text_frame:
            continue

        original = shape.text

        # Substituições simples
        texto = original.replace("{DIA}", row["DIA"])
        texto = texto.replace("{MES}", row["MES"])
        texto = texto.replace("{ANO}", row["ANO"])
        texto = texto.replace("{HORAS}", str(row["Horas"]))

        # Caixa de texto no slide novo
        textbox = slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
        tf = textbox.text_frame
        tf.word_wrap = True
        p = tf.add_paragraph()

        i = 0
        while i < len(texto):
            achou = False
            for ph, conf in placeholders_azuis.items():
                if texto[i:].startswith(ph):
                    run = p.add_run()
                    run.text = str(row[conf["valor"]])
                    run.font.name = "Arial Rounded MT Bold"
                    run.font.size = conf["size"]
                    run.font.color.rgb = RGBColor(0, 112, 192)  # azul
                    i += len(ph)
                    achou = True
                    break
            if not achou:
                run = p.add_run()
                run.text = texto[i]
                run.font.name = "Arial Rounded MT Bold"
                run.font.size = Pt(16)
                run.font.color.rgb = RGBColor(0, 0, 0)  # preto
                i += 1

    # Remove o slide do template (primeiro slide)
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[0])

    # Nome do arquivo
    nome_base = f"certificado_{row['Nome']}_{row['Capacitacao']}"
    nome_base = re.sub(r'[<>:"/\\|?*]', '_', nome_base)

    pptx_path = os.path.join("certificados", nome_base + ".pptx")
    prs.save(pptx_path)
    print(f"✅ PPTX gerado sem slide extra: {pptx_path}")
